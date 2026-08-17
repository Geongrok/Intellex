"""Knowledge Base module.

Ingests documents (PDF, DOCX, PPTX, XLSX, TXT/MD/CSV, and images via OCR),
splits them into chunks, and builds BOTH:

  * a BM25 keyword index (fast, no deps), and
  * a ChromaDB vector index (semantic, via auto-detected embeddings)

so the chatbot can retrieve the most relevant passages using hybrid
(vector + keyword) search.
"""

import hashlib
import os
import pickle
import re
from typing import Dict, List, Optional, Tuple

# The project root is the directory that contains the "app" package.
# Anchoring data/cache here makes the chatbot work regardless of the
# current working directory the server is launched from.
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_DATA_DIR = os.path.join(PROJECT_ROOT, "data")
DEFAULT_CACHE_DIR = os.path.join(PROJECT_ROOT, "cache")

from rank_bm25 import BM25Okapi

from .vector_store import VectorStore

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    import docx  # python-docx
except ImportError:  # pragma: no cover
    docx = None

try:
    from pptx import Presentation  # python-pptx
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import openpyxl
except ImportError:  # pragma: no cover
    openpyxl = None

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

try:
    from rapidocr_onnxruntime import RapidOCR

    _OCR = RapidOCR()
    _OCR_AVAILABLE = True
except Exception:  # OCR optional
    _OCR = None
    _OCR_AVAILABLE = False

SUPPORTED_EXTS = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".csv",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp",
}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp"}
CHUNK_SIZE = 900
CHUNK_OVERLAP = 120

# Common English stop words excluded from relevance scoring so generic
# questions like "What is the capital of France?" don't falsely match a
# document about engineering just because of "what/is/the/of".
STOP_WORDS = set(
    """
    a an and are as at be but by for from have has he her his i if in is it
    its of on or our she so that the their them then there these they this
    to was we what when which who will with you your
    """.split()
)


class KnowledgeBase:
    """Builds and queries BM25 + vector indexes over user documents."""

    def __init__(
        self,
        data_dir: Optional[str] = None,
        cache_dir: Optional[str] = None,
    ):
        # Default to the project-root-anchored "data"/"cache" folders so the
        # chatbot finds your documents no matter where the server is started.
        self.data_dir = data_dir or DEFAULT_DATA_DIR
        self.cache_dir = cache_dir or DEFAULT_CACHE_DIR
        self.chunks: List[str] = []
        self.meta: List[Dict] = []  # {"file": str, "page": int|None}
        self.tokenized: List[List[str]] = []
        self.bm25: Optional[BM25Okapi] = None
        self.vector = VectorStore(
            persist_dir=os.path.join(self.cache_dir, "chroma")
        )
        os.makedirs(self.data_dir, exist_ok=True)
        os.makedirs(self.cache_dir, exist_ok=True)
        self._cache_path = os.path.join(self.cache_dir, "index.pkl")

    # ------------------------------------------------------------------ #
    # Public API
    # ------------------------------------------------------------------ #
    def build_index(self, force: bool = False) -> None:
        """Load the cached index if valid, otherwise (re)build it."""
        if not force and os.path.exists(self._cache_path):
            try:
                with open(self._cache_path, "rb") as fh:
                    cached = pickle.load(fh)
                if cached.get("signature") == self._signature():
                    self.chunks = cached["chunks"]
                    self.meta = cached["meta"]
                    self._build_bm25()
                    # If the cached index somehow has zero documents but the
                    # data folder is populated, rebuild to be safe.
                    if not self.chunks and self._list_files():
                        self.rebuild()
                    return
            except Exception:
                pass
        self.rebuild()

    def rebuild(self) -> None:
        """Re-parse every document and re-index from scratch."""
        self.chunks = []
        self.meta = []
        all_ids = []
        all_texts = []
        all_meta = []

        for path in self._list_files():
            try:
                segments = self._extract_segments(path)
            except Exception as exc:  # skip unreadable files
                print(f"[KB] Skipped {os.path.basename(path)}: {exc}")
                continue
            for seg_idx, (text, page) in enumerate(segments):
                for piece_idx, piece in enumerate(self._chunk(text)):
                    if not piece:
                        continue
                    self.chunks.append(piece)
                    self.meta.append(
                        {"file": os.path.basename(path), "page": page}
                    )
                    # Unique ID: full path + page + segment + chunk index + content hash,
                    # so identical text from different files/folders never collides.
                    rel = os.path.relpath(path, self.data_dir)
                    uid = hashlib.md5(
                        f"{rel}|{page}|{seg_idx}|{piece_idx}|{piece}".encode("utf-8")
                    ).hexdigest()
                    all_ids.append(uid)
                    all_texts.append(piece)
                    all_meta.append(
                        {"file": os.path.basename(path), "page": page or ""}
                    )

        self._build_bm25()

        # Rebuild the vector store.
        self.vector.clear()
        if all_texts:
            self.vector.add_documents(all_ids, all_texts, all_meta)

        self._save_cache()

    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        """Hybrid search: fuse BM25 (keyword) + vector (semantic) scores.

        Returns top-k chunks with a 0..1 relevance score.
        """
        keyword_results = self._search_keyword(query, top_k=top_k)
        vector_results = self.vector.query(query, top_k=top_k)

        # Fuse by text (chunk body) so duplicates from both retrievers merge.
        fused: Dict[str, Dict] = {}
        for r in keyword_results:
            key = r["text"][:200]
            fused[key] = dict(r)
        for r in vector_results:
            key = r["text"][:200]
            if key in fused:
                # Merge: take the max score, keep keyword's file/page.
                fused[key]["score"] = max(fused[key]["score"], r["score"])
                fused[key]["vector_score"] = r["score"]
            else:
                fused[key] = {
                    "text": r["text"],
                    "file": r["file"],
                    "page": r["page"],
                    "score": r["score"] * 0.85,  # slight discount vector-only
                    "vector_score": r["score"],
                }

        results = sorted(fused.values(), key=lambda x: x["score"], reverse=True)
        # Filter out anything with effectively zero relevance.
        results = [r for r in results if r["score"] > 0]
        return results[:top_k]

    def _search_keyword(self, query: str, top_k: int = 5) -> List[Dict]:
        """Pure BM25 keyword search (used as part of hybrid search).

        Relevance is computed as the fraction of the query's meaningful
        (non-stop-word) terms that appear in the chunk.  Chunks sharing only
        a single meaningful term are heavily penalized so unrelated questions
        don't accidentally rank as a "database hit".
        """
        if not self.bm25 or not self.chunks:
            return []
        query_tokens = self.tokenize(query)
        if not query_tokens:
            return []
        unique_query = set(t for t in query_tokens if t not in STOP_WORDS)
        if not unique_query:
            unique_query = set(query_tokens)
        scores = self.bm25.get_scores(query_tokens)
        ranked = sorted(
            zip(scores, range(len(self.chunks))),
            key=lambda x: x[0],
            reverse=True,
        )

        # Pre-compute filename tokens once for year/filename boosting.
        # e.g. "2021.pdf" -> {"2021"}; "Math.docx" -> {"math"}.
        file_tokens_cache = [
            self.tokenize(self.meta[idx]["file"]) for idx in range(len(self.chunks))
        ]
        # Scan a wide window (up to all chunks) so year/filename matches are
        # not cut off by BM25's raw-frequency ranking.
        window = max(top_k * 6, 200)
        window = min(window, len(self.chunks))

        results = []
        for _bm25, idx in ranked[:window]:
            chunk_tokens = set(self.tokenized[idx])
            hits = unique_query & chunk_tokens
            coverage = len(hits) / len(unique_query)
            if coverage <= 0:
                continue
            # Penalize single-term matches: one coincidental word is not a
            # genuine topical match.
            if len(hits) == 1:
                coverage -= 0.35

            score = max(float(coverage), 0.0)

            # Year/filename boost: if the query contains a token (usually a
            # year like "2021") that appears in the source filename, strongly
            # promote that file. This makes "GATE 2021 papers" correctly pick
            # 2021.pdf instead of a generic 2024/2014 header.
            if any(t in file_tokens_cache[idx] for t in unique_query):
                score = min(1.0, score + 0.45)

            results.append(
                {
                    "text": self.chunks[idx],
                    "file": self.meta[idx]["file"],
                    "page": self.meta[idx]["page"],
                    "score": score,  # 0..1 relevance
                    "bm25": float(_bm25),
                    "filename_match": any(
                        t in file_tokens_cache[idx] for t in unique_query
                    ),
                }
            )

        # Re-rank by the boosted relevance score (not raw BM25).
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:top_k]

    def doc_count(self) -> int:
        return len(self.chunks)

    def file_names(self) -> List[str]:
        names = set()
        for m in self.meta:
            names.add(m["file"])
        return sorted(names)

    # ------------------------------------------------------------------ #
    # Index helpers
    # ------------------------------------------------------------------ #
    @staticmethod
    def tokenize(text: str) -> List[str]:
        return re.findall(r"\w+", text.lower())

    def _chunk_id(self, text: str) -> str:
        return hashlib.md5(text.encode("utf-8")).hexdigest()

    def _build_bm25(self) -> None:
        self.tokenized = [self.tokenize(c) for c in self.chunks]
        self.bm25 = BM25Okapi(self.tokenized) if self.tokenized else None

    def _signature(self) -> str:
        parts = []
        for path in self._list_files():
            st = os.stat(path)
            parts.append(f"{path}:{st.st_size}:{int(st.st_mtime)}")
        return hashlib.md5("\n".join(parts).encode("utf-8")).hexdigest()

    def _save_cache(self) -> None:
        try:
            with open(self._cache_path, "wb") as fh:
                pickle.dump(
                    {
                        "signature": self._signature(),
                        "chunks": self.chunks,
                        "meta": self.meta,
                    },
                    fh,
                )
        except Exception as exc:  # cache is optional
            print(f"[KB] Could not save cache: {exc}")

    def _list_files(self) -> List[str]:
        found = []
        if not os.path.isdir(self.data_dir):
            return found
        for root, _dirs, files in os.walk(self.data_dir):
            for fname in sorted(files):
                if fname.startswith("~$"):  # Office temp lock files
                    continue
                ext = os.path.splitext(fname)[1].lower()
                if ext in SUPPORTED_EXTS:
                    found.append(os.path.join(root, fname))
        return found

    # ------------------------------------------------------------------ #
    # Text extraction
    # ------------------------------------------------------------------ #
    def _extract_segments(self, path: str) -> List[Tuple[str, Optional[int]]]:
        """Return a list of (text, page_or_slide_number) segments."""
        ext = os.path.splitext(path)[1].lower()

        if ext == ".pdf":
            return self._extract_pdf(path)
        if ext == ".docx":
            return self._extract_docx(path)
        if ext == ".pptx":
            return self._extract_pptx(path)
        if ext == ".xlsx":
            return self._extract_xlsx(path)
        if ext in IMAGE_EXTS:
            return self._extract_image(path)
        if ext in (".txt", ".md", ".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                return [(fh.read(), None)]
        return []

    def _extract_pdf(self, path: str) -> List[Tuple[str, Optional[int]]]:
        if PdfReader is None:
            return [("[PDF support not installed]", None)]
        reader = PdfReader(path)
        segments = []
        for i, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            if text.strip():
                segments.append((text, i))
        return segments

    def _extract_docx(self, path: str) -> List[Tuple[str, Optional[int]]]:
        if docx is None:
            return [("[DOCX support not installed]", None)]
        d = docx.Document(path)
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    parts.append(line)
        return [("\n".join(parts), None)]

    def _extract_pptx(self, path: str) -> List[Tuple[str, Optional[int]]]:
        if Presentation is None:
            return [("[PPTX support not installed]", None)]
        prs = Presentation(path)
        segments = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            texts.append(t)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        line = " | ".join(c for c in cells if c)
                        if line:
                            texts.append(line)
            if texts:
                segments.append(("\n".join(texts), i))
        return segments

    def _extract_xlsx(self, path: str) -> List[Tuple[str, Optional[int]]]:
        if openpyxl is None:
            return [("[XLSX support not installed]", None)]
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        segments = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if vals:
                    rows.append(" | ".join(vals))
            if rows:
                segments.append(
                    (f"[Sheet: {ws.title}]\n" + "\n".join(rows), None)
                )
        return segments

    def _extract_image(self, path: str) -> List[Tuple[str, Optional[int]]]:
        """OCR an image file into text using RapidOCR (self-contained ONNX)."""
        if Image is None or not _OCR_AVAILABLE:
            return [("[OCR not available for this image]", None)]
        try:
            img = Image.open(path).convert("RGB")
            # RapidOCR expects a numpy array or file path.
            result, _ = _OCR(str(path))
            if not result:
                return []
            text = "\n".join(line[1] for line in result if len(line) > 1)
            if text.strip():
                return [(f"[Image OCR: {os.path.basename(path)}]\n{text}", None)]
            return []
        except Exception as exc:
            return [(f"[Image OCR failed: {exc}]", None)]

    # ------------------------------------------------------------------ #
    # Chunking
    # ------------------------------------------------------------------ #
    @staticmethod
    def _chunk(text: str) -> List[str]:
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            return []
        if len(text) <= CHUNK_SIZE:
            return [text]
        chunks = []
        start = 0
        while start < len(text):
            end = start + CHUNK_SIZE
            if end < len(text):
                cut = text.rfind(" ", start, end)
                if cut > start + CHUNK_SIZE // 2:
                    end = cut
            chunks.append(text[start:end].strip())
            start = max(end - CHUNK_OVERLAP, start + 1)
        return chunks
